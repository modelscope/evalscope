from __future__ import annotations

import base64
import binascii
import hashlib
import json
import re
from contextlib import closing
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple
from zipfile import BadZipFile, ZipFile

from evalscope.api.agent import AgentEnvironment
from evalscope.api.agent.types import ExecResult
from evalscope.api.messages import ChatMessageSystem, ChatMessageUser, ContentImage, ContentText
from evalscope.utils.io_utils import safe_filename

SANDBOX_REFERENCE_DIR = 'reference_files'
SANDBOX_OUTPUT_DIR = 'jobbench_output'
HOST_ARTIFACT_ROOT = 'artifacts/job_bench'

MAX_CHARS_PER_FILE_DEFAULT = 200_000
MAX_VISION_IMAGES_DEFAULT = 8
SQLITE_EXTS = {'db', 'sqlite', 'sqlite3'}
SQLITE_ROWS_PER_TABLE = 500
VISION_IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
VISION_MIME = {
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'gif': 'image/gif',
    'webp': 'image/webp',
}
VISUAL_RUBRIC_PATTERN = re.compile(
    r'\b(plot|figure|visualization|visualisation|visualize|visualise|'
    r'heatmap|histogram|scatter ?plot|biplot|diagram|q[- ]?q)\b',
    re.IGNORECASE,
)


class JobBenchArtifactEnvironment(AgentEnvironment):
    name = 'job_bench_artifact'

    def __init__(self, env: AgentEnvironment, artifact_dir: Path, metadata: Dict[str, Any]) -> None:
        self._env = env
        self._artifact_dir = artifact_dir
        self._metadata = metadata

    async def exec(
        self,
        cmd: List[str],
        *,
        cwd: Optional[str] = None,
        input: Optional[str] = None,
        timeout: Optional[float] = None,
        env: Optional[Dict[str, str]] = None,
    ) -> ExecResult:
        return await self._env.exec(cmd, cwd=cwd, input=input, timeout=timeout, env=env)

    async def close(self) -> None:
        try:
            await self._env.close()
        finally:
            output_dir = self._artifact_dir / SANDBOX_OUTPUT_DIR
            self._metadata['artifact_dir'] = str(self._artifact_dir)
            self._metadata['output_files'] = [{
                'path': f'{SANDBOX_OUTPUT_DIR}/{path.relative_to(output_dir).as_posix()}',
                'local_path': str(path),
            } for path in sorted(output_dir.rglob('*')) if path.is_file()]


def artifact_dir(sample: Any, output_dir: str) -> Path:
    task_id = str(sample.metadata.get('task_id') or sample.id or 'unknown')
    return (Path(output_dir) / HOST_ARTIFACT_ROOT / safe_filename(task_id)).resolve()


def parse_rubrics(value: Any) -> List[Dict[str, Any]]:
    if value is None:
        return []
    if isinstance(value, str):
        try:
            value = json.loads(value)
        except json.JSONDecodeError:
            return []
    if isinstance(value, list):
        return [item for item in value if isinstance(item, dict)]
    if isinstance(value, dict):
        rubrics = value.get('rubrics') or value.get('evaluation_rubrics') or []
        return [item for item in rubrics if isinstance(item, dict)]
    return []


def normalize_criteria(rubric: Dict[str, Any]) -> List[str]:
    criterion_raw = rubric.get('criterion', [])
    if isinstance(criterion_raw, str):
        return [criterion_raw]
    if isinstance(criterion_raw, list):
        return [str(item) for item in criterion_raw]
    return []


def build_failed_rubric_result(
    rubric_index: int,
    rubric: Dict[str, Any],
    overall_reasoning: str,
) -> Dict[str, Any]:
    criteria = normalize_criteria(rubric)
    return {
        'index': rubric_index,
        'rubric': rubric.get('rubric', ''),
        'weight': rubric.get('weight', 0),
        'result': {
            'passed': False,
            'score': 0,
            'criteria_count': len(criteria),
            'criteria_passed': 0,
            'criteria_results': [{
                'index': idx,
                'criterion': criterion,
                'passed': False,
                'reasoning': overall_reasoning,
                'evidence': '',
            } for idx, criterion in enumerate(criteria)],
            'overall_reasoning': overall_reasoning,
        },
    }


def build_scorecard(results: List[Dict[str, Any]]) -> Dict[str, Any]:
    total_score = sum(float(result['result']['score']) for result in results)
    max_score = sum(float(result.get('weight') or 0) for result in results)
    passed_count = sum(1 for result in results if result['result']['passed'])
    total_count = len(results)
    normalized = round(total_score / max_score, 4) if max_score > 0 else 0.0
    pass_rate = round(passed_count / total_count, 4) if total_count > 0 else 0.0
    return {
        'total_score': total_score,
        'max_score': max_score,
        'normalized_score': normalized,
        'pass_rate': pass_rate,
        'passed_count': passed_count,
        'total_count': total_count,
    }


def convert_file_to_text(path: Path, max_chars: int = MAX_CHARS_PER_FILE_DEFAULT) -> str:
    ext = path.suffix.lower().lstrip('.')
    if ext in (
        'txt',
        'md',
        'csv',
        'py',
        'json',
        'sh',
        'log',
        'xml',
        'html',
        'css',
        'js',
        'ts',
        'yaml',
        'yml',
        'ini',
        'cfg',
        'conf',
        'sql',
        'rules',
        'geojson',
    ):
        content = _read_text(path)
    elif ext in ('xlsx', 'xls'):
        content = _excel_to_text(path)
    elif ext == 'docx':
        content = _docx_to_text(path)
    elif ext == 'pdf':
        content = _pdf_to_text(path)
    elif ext in SQLITE_EXTS:
        content = _sqlite_to_text(path)
    elif ext == 'pptx':
        content = _pptx_to_text(path)
    elif ext == 'ipynb':
        content = _notebook_to_text(path)
    elif ext in ('png', 'jpg', 'jpeg', 'gif', 'svg', 'bmp', 'webp'):
        content = f'[Image file: {path.name} - cannot extract text content]'
    else:
        content = f'[Binary or unsupported file type: {ext} - {path.name}]'
    if ext not in SQLITE_EXTS and len(content) > max_chars:
        return content[:max_chars] + f'\n... [Content truncated at {max_chars} characters]'
    return content


def extract_all_file_contents(output_dir: Path, max_chars_per_file: int = MAX_CHARS_PER_FILE_DEFAULT) -> str:
    parts = []
    for file_path in sorted(output_dir.rglob('*')):
        if not file_path.is_file():
            continue
        content = convert_file_to_text(file_path, max_chars=max_chars_per_file)
        relative_name = file_path.relative_to(output_dir).as_posix()
        parts.append(f'=== FILE: {relative_name} ===\n{content}\n')
    return '\n'.join(parts)


def rubric_needs_vision(rubric: Dict[str, Any]) -> bool:
    text = rubric.get('rubric', '') or ''
    criterion = rubric.get('criterion', [])
    if isinstance(criterion, list):
        text = text + ' ' + ' '.join(str(item) for item in criterion)
    elif isinstance(criterion, str):
        text = text + ' ' + criterion
    return bool(VISUAL_RUBRIC_PATTERN.search(text))


def collect_image_attachments(output_dir: Path, cap: int = MAX_VISION_IMAGES_DEFAULT) -> List[Tuple[str, str]]:
    if not output_dir.exists() or cap <= 0:
        return []

    attachments: List[Tuple[str, str]] = []
    seen_hashes = set()

    def add_attachment(name: str, mime: str, image_bytes: bytes) -> bool:
        digest = hashlib.sha256(image_bytes).hexdigest()
        if digest in seen_hashes:
            return False
        seen_hashes.add(digest)
        encoded = base64.b64encode(image_bytes).decode('ascii')
        attachments.append((name, f'data:{mime};base64,{encoded}'))
        return len(attachments) >= cap

    for path in sorted(output_dir.rglob('*')):
        if not path.is_file() or path.suffix.lower().lstrip('.') not in VISION_IMAGE_EXTS:
            continue
        try:
            image_bytes = path.read_bytes()
        except OSError:
            continue
        ext = path.suffix.lower().lstrip('.')
        if add_attachment(path.relative_to(output_dir).as_posix(), VISION_MIME[ext], image_bytes):
            return attachments

    for docx_path in sorted(output_dir.rglob('*.docx')):
        try:
            with ZipFile(docx_path) as archive:
                for media_name in sorted(archive.namelist()):
                    ext = Path(media_name).suffix.lower().lstrip('.')
                    if not media_name.startswith('word/media/') or ext not in VISION_IMAGE_EXTS:
                        continue
                    name = f'{docx_path.relative_to(output_dir).as_posix()}:{media_name}'
                    if add_attachment(name, VISION_MIME[ext], archive.read(media_name)):
                        return attachments
        except (BadZipFile, OSError):
            continue

    for notebook_path in sorted(output_dir.rglob('*.ipynb')):
        try:
            notebook = json.loads(notebook_path.read_text(encoding='utf-8'))
        except (OSError, json.JSONDecodeError):
            continue
        for cell_index, cell in enumerate(notebook.get('cells', [])):
            for output_index, output in enumerate(cell.get('outputs', [])):
                data = output.get('data', {})
                if not isinstance(data, dict):
                    continue
                for mime in VISION_MIME.values():
                    encoded = data.get(mime)
                    if isinstance(encoded, list):
                        encoded = ''.join(str(part) for part in encoded)
                    if not isinstance(encoded, str):
                        continue
                    try:
                        image_bytes = base64.b64decode(''.join(encoded.split()), validate=True)
                    except (ValueError, binascii.Error):
                        continue
                    name = (
                        f'{notebook_path.relative_to(output_dir).as_posix()}:'
                        f'cell-{cell_index}-output-{output_index}'
                    )
                    if add_attachment(name, mime, image_bytes):
                        return attachments
    return attachments


def build_judge_prompt(rubric: Dict[str, Any], file_contents: str, vision_used: bool) -> str:
    rubric_text = rubric.get('rubric', '')
    criteria = normalize_criteria(rubric)
    criteria_list_text = '\n'.join(f'Criterion {idx}: {criterion}' for idx, criterion in enumerate(criteria))
    return f"""You are an evaluation judge. Your task is to evaluate ALL criteria for a single rubric.

## Rubric Description
{rubric_text}

## Criteria to Evaluate (Judge ALL of them)
{criteria_list_text}

## Output Files Content
The following are the contents of all output files to evaluate:

{file_contents}

## Evaluation Rules
- Evaluate EACH criterion listed above independently
- For each criterion: determine if it PASSES or FAILS
- Semantic matching is acceptable (you don't need exact wording match)
- Binary judgment for each criterion: PASS or FAIL only
- The rubric passes ONLY if ALL criteria pass

## Output Format
Return your judgment as a JSON object with EXACTLY this structure (no markdown, no extra text):
{{
  "criteria_results": [
    {{"index": 0, "passed": true/false, "reasoning": "...", "evidence": "..."}},
    {{"index": 1, "passed": true/false, "reasoning": "...", "evidence": "..."}}
  ],
  "rubric_passed": true/false,
  "overall_reasoning": "Summary of why the rubric passed or failed"
}}

IMPORTANT:
- criteria_results array must have exactly {len(criteria)} items (one for each criterion)
- rubric_passed should be true ONLY if ALL criteria passed
- Include specific evidence from the output files{' and the attached images' if vision_used else ''}
"""


def parse_judge_json(content: str) -> Dict[str, Any]:
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        pass

    fence = re.search(r'```(?:json)?\s*\n(.*?)\n\s*```', content, re.DOTALL)
    if fence:
        try:
            return json.loads(fence.group(1).strip())
        except json.JSONDecodeError:
            pass

    first = content.find('{')
    last = content.rfind('}')
    if first != -1 and last > first:
        try:
            return json.loads(content[first:last + 1])
        except json.JSONDecodeError:
            pass
    raise ValueError(f'Could not extract JSON from response: {content[:500]}')


def judge_rubric(
    *,
    rubric_index: int,
    rubric: Dict[str, Any],
    file_contents: str,
    judge: Callable[..., str],
    image_attachments: Optional[List[Tuple[str, str]]] = None,
) -> Dict[str, Any]:
    weight = rubric.get('weight', 0)
    criteria = normalize_criteria(rubric)
    attached_images = list(image_attachments or []) if rubric_needs_vision(rubric) else []
    vision_used = bool(attached_images)
    prompt = build_judge_prompt(rubric, file_contents, vision_used=vision_used)

    try:
        if vision_used:
            content = [ContentText(text=prompt)]
            content.append(ContentText(text=f'\n## Attached Images ({len(attached_images)} files)'))
            for idx, (name, url) in enumerate(attached_images, start=1):
                content.append(ContentText(text=f'Image {idx}: {name}'))
                content.append(ContentImage(image=url))
            raw_response = judge(
                messages=[
                    ChatMessageSystem(content='You are an evaluation judge. Return valid JSON only, with no markdown.'),
                    ChatMessageUser(content=content),
                ]
            )
        else:
            raw_response = judge(
                prompt=prompt,
                system_prompt='You are an evaluation judge. Return valid JSON only, with no markdown.',
            )
        parsed = parse_judge_json(raw_response.strip())
    except Exception as exc:
        return build_failed_rubric_result(rubric_index, rubric, f'Judge failed: {exc}')

    model_criteria = parsed.get('criteria_results', [])
    criteria_results = []
    for idx, criterion in enumerate(criteria):
        item = model_criteria[idx] if idx < len(model_criteria) and isinstance(model_criteria[idx], dict) else {}
        criteria_results.append({
            'index': idx,
            'criterion': criterion,
            'passed': bool(item.get('passed', False)),
            'reasoning': item.get('reasoning', ''),
            'evidence': item.get('evidence', ''),
        })

    rubric_passed = bool(parsed.get('rubric_passed', False))
    return {
        'index': rubric_index,
        'rubric': rubric.get('rubric', ''),
        'weight': weight,
        'result': {
            'passed': rubric_passed,
            'score': weight if rubric_passed else 0,
            'criteria_count': len(criteria),
            'criteria_passed': sum(1 for item in criteria_results if item['passed']),
            'criteria_results': criteria_results,
            'overall_reasoning': parsed.get('overall_reasoning', ''),
        },
    }


def evaluate_job_bench_output(
    *,
    output_dir: Path,
    rubrics: List[Dict[str, Any]],
    judge: Callable[..., str],
) -> Dict[str, Any]:
    if not rubrics:
        return {
            'scorecard': build_scorecard([]),
            'rubrics': [],
            'error': 'no rubrics',
        }
    if not output_dir.exists() or not any(path.is_file() for path in output_dir.rglob('*')):
        results = [
            build_failed_rubric_result(idx, rubric, 'No output files found in the model output directory.')
            for idx, rubric in enumerate(rubrics)
        ]
        return {'scorecard': build_scorecard(results), 'rubrics': results}

    file_contents = extract_all_file_contents(output_dir)
    if not file_contents.strip():
        results = [
            build_failed_rubric_result(idx, rubric, 'Output files were unreadable or empty after conversion.')
            for idx, rubric in enumerate(rubrics)
        ]
        return {'scorecard': build_scorecard(results), 'rubrics': results}

    image_attachments = collect_image_attachments(output_dir)
    results = []
    for idx, rubric in enumerate(rubrics):
        results.append(
            judge_rubric(
                rubric_index=idx,
                rubric=rubric,
                file_contents=file_contents,
                judge=judge,
                image_attachments=image_attachments,
            )
        )
    return {
        'scorecard': build_scorecard(results),
        'rubrics': results,
    }


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding='utf-8', errors='replace')
    except Exception as exc:
        return f'[ERROR: Failed to read text file: {path.name}: {exc}]'


def _excel_to_text(path: Path) -> str:
    try:
        import pandas as pd

        xl = pd.ExcelFile(str(path))
        parts = []
        for sheet in xl.sheet_names:
            df = pd.read_excel(xl, sheet_name=sheet)
            parts.append(f'=== Sheet: {sheet} ===\n{df.to_csv(index=False)}')
        return '\n'.join(parts)
    except ImportError:
        return f'[ERROR: pandas/openpyxl not available for {path.name}]'
    except Exception as exc:
        return f'[ERROR: Failed to read Excel {path.name}: {exc}]'


def _docx_to_text(path: Path) -> str:
    try:
        import mammoth

        with open(str(path), 'rb') as f:
            result = mammoth.convert_to_markdown(f)
        return result.value
    except ImportError:
        return f'[ERROR: mammoth not available for {path.name}]'
    except Exception as exc:
        return f'[ERROR: Failed to read DOCX {path.name}: {exc}]'


def _pdf_to_text(path: Path) -> str:
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            parts = []
            for idx, page in enumerate(pdf.pages):
                text = page.extract_text(layout=True) or ''
                parts.append(f'=== Page {idx + 1} ===\n{text}')
        return '\n'.join(parts)
    except ImportError:
        return f'[ERROR: pdfplumber not available for {path.name}]'
    except Exception as exc:
        return f'[ERROR: Failed to read PDF {path.name}: {exc}]'


def _sqlite_to_text(path: Path) -> str:
    import sqlite3 as sqlite

    try:
        with closing(sqlite.connect(str(path))) as con:
            cur = con.cursor()
            schema = con.execute('SELECT sql FROM sqlite_master WHERE sql IS NOT NULL').fetchall()
            parts = ['=== Schema ===']
            parts.extend(row[0] for row in schema if row[0])
            tables = [row[0] for row in con.execute("SELECT name FROM sqlite_master WHERE type='table'").fetchall()]
            for table in tables:
                parts.append(f'\n=== Table: {table} ===')
                try:
                    total_rows = con.execute(f'SELECT COUNT(*) FROM "{table}"').fetchone()[0]
                    rows = cur.execute(f'SELECT * FROM "{table}" LIMIT {SQLITE_ROWS_PER_TABLE}').fetchall()
                    cols = [d[0] for d in cur.description]
                    parts.append(f'-- total_rows: {total_rows}; shown: {len(rows)} (LIMIT {SQLITE_ROWS_PER_TABLE})')
                    parts.append(','.join(cols))
                    for row in rows:
                        parts.append(','.join('' if value is None else str(value) for value in row))
                except Exception as exc:
                    parts.append(f'[ERROR reading table {table}: {exc}]')
        return '\n'.join(parts)
    except Exception as exc:
        return f'[ERROR: Failed to read SQLite {path.name}: {exc}]'


def _pptx_to_text(path: Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []
        for idx, slide in enumerate(prs.slides):
            parts.append(f'=== Slide {idx + 1} ===')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    parts.append(shape.text)
        return '\n'.join(parts)
    except ImportError:
        return f'[ERROR: python-pptx not available for {path.name}]'
    except Exception as exc:
        return f'[ERROR: Failed to read PowerPoint {path.name}: {exc}]'


def _notebook_to_text(path: Path) -> str:
    try:
        nb = json.loads(path.read_text(encoding='utf-8'))
        parts = []
        for cell in nb.get('cells', []):
            parts.append(f'=== {cell.get("cell_type", "unknown")} ===')
            parts.append(''.join(cell.get('source', [])))
            for output in cell.get('outputs', []):
                if 'text' in output:
                    parts.append(''.join(output['text']))
        return '\n'.join(parts)
    except Exception as exc:
        return f'[ERROR: Failed to read notebook {path.name}: {exc}]'
